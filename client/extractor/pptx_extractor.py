from pptx import Presentation
from io import BytesIO
from utils.logging import get_logger

logger = get_logger(__name__)


def extract_pptx(file_bytes: bytes) -> str:
    """
    Extracts text from a PPTX file (bytes).

    - Loads presentation from memory
    - Iterates through slides and shapes
    - Extracts text only from shapes with text frames
    - Prefixes each slide with: --- Slide N ---
    - Skips slides with no text
    """

    try:
        presentation = Presentation(BytesIO(file_bytes))

        slides_text = []

        for i, slide in enumerate(presentation.slides, start=1):
            slide_lines = []

            for shape in slide.shapes:
                if not hasattr(shape, "text_frame") or not shape.text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_lines.append(text)

            if not slide_lines:
                continue  # skip empty slides

            slide_block = f"--- Slide {i} ---\n" + "\n".join(slide_lines)
            slides_text.append(slide_block)

        if not slides_text:
            logger.warning("PPTX extraction completed but no text found.")
            return ""

        return "\n\n".join(slides_text)

    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return ""