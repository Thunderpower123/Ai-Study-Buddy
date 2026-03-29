# test_ingest.py
# Tests the full ingestion pipeline layer by layer.
#
# HOW TO RUN (from inside /client with venv active):
#   python -m pytest tests/test_ingest.py -v
#
# What gets tested:
#   1. PDF extractor  — can it pull text out of a real PDF?
#   2. PPTX extractor — can it pull text out of a real PPTX?
#   3. DOCX extractor — can it pull text out of a real DOCX?
#   4. Chunker        — does it split correctly with overlap?
#   5. Embeddings     — does OpenAI return a real 1536-dim vector?
#   6. Pinecone upsert + search — do vectors actually go in and come back?
#   7. Full /ingest HTTP endpoint — end-to-end with a real PDF over HTTP

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import base64
import asyncio
import pytest
import httpx

# ── Minimal in-memory PDF (real PDF bytes, not a mock) ──────────────────────
# This is a tiny but valid 1-page PDF containing the text "Newton second law"
# Generated with reportlab so you don't need a file on disk.
def make_test_pdf_bytes() -> bytes:
    from reportlab.pdfgen import canvas
    from io import BytesIO
    buf = BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(100, 750, "Newton second law states F equals ma.")
    c.drawString(100, 730, "Force is mass times acceleration.")
    c.drawString(100, 710, "This is a fundamental principle of classical mechanics.")
    c.save()
    return buf.getvalue()


def make_test_docx_bytes() -> bytes:
    from docx import Document
    from io import BytesIO
    doc = Document()
    doc.add_paragraph("Fleming right hand rule determines motor direction.")
    doc.add_paragraph("The thumb points in the direction of motion.")
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_test_pptx_bytes() -> bytes:
    from pptx import Presentation
    from io import BytesIO
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Ohm Law"
    slide.placeholders[1].text = "Voltage equals current times resistance."
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── 1. PDF Extractor ─────────────────────────────────────────────────────────
def test_pdf_extractor_returns_text():
    from extractor.pdf_extractor import extract_pdf
    pdf_bytes = make_test_pdf_bytes()
    text = extract_pdf(pdf_bytes)
    assert isinstance(text, str), "Should return a string"
    assert len(text) > 10, "Should return non-empty text"
    assert "Newton" in text or "Force" in text or "mass" in text, (
        f"Expected physics content in extracted text, got: {text[:200]}"
    )
    print(f"\n[PDF] Extracted {len(text)} chars. Preview: {text[:100]}")


# ── 2. DOCX Extractor ────────────────────────────────────────────────────────
def test_docx_extractor_returns_text():
    from extractor.docx_extractor import extract_docx
    docx_bytes = make_test_docx_bytes()
    text = extract_docx(docx_bytes)
    assert isinstance(text, str)
    assert "Fleming" in text or "motor" in text or "thumb" in text, (
        f"Expected Fleming content, got: {text[:200]}"
    )
    print(f"\n[DOCX] Extracted {len(text)} chars. Preview: {text[:100]}")


# ── 3. PPTX Extractor ────────────────────────────────────────────────────────
def test_pptx_extractor_returns_text():
    from extractor.pptx_extractor import extract_pptx
    pptx_bytes = make_test_pptx_bytes()
    text = extract_pptx(pptx_bytes)
    assert isinstance(text, str)
    assert "Ohm" in text or "Voltage" in text or "resistance" in text, (
        f"Expected Ohm content, got: {text[:200]}"
    )
    print(f"\n[PPTX] Extracted {len(text)} chars. Preview: {text[:100]}")


# ── 4. Chunker ───────────────────────────────────────────────────────────────
def test_chunker_basic():
    from extractor.chunker import chunk_text
    # 600 words → should produce 2 chunks (500 words chunk1, overlap 50, chunk2 starts at 450)
    text = " ".join([f"word{i}" for i in range(600)])
    chunks = chunk_text(text, chunk_size=500, overlap=50)
    assert len(chunks) >= 2, f"Expected >=2 chunks, got {len(chunks)}"
    print(f"\n[CHUNKER] {len(chunks)} chunks from 600 words")


def test_chunker_overlap():
    from extractor.chunker import chunk_text
    words = [f"w{i}" for i in range(600)]
    text = " ".join(words)
    chunks = chunk_text(text, chunk_size=500, overlap=50)
    # Last 50 words of chunk 1 should appear at start of chunk 2
    chunk1_last = chunks[0].split()[-50:]
    chunk2_first = chunks[1].split()[:50]
    assert chunk1_last == chunk2_first, "Overlap words should match between chunks"
    print(f"\n[CHUNKER] Overlap verified: last 50 of chunk1 == first 50 of chunk2")


def test_chunker_empty_input():
    from extractor.chunker import chunk_text
    chunks = chunk_text("")
    assert chunks == [], "Empty input should return empty list"


def test_chunker_short_text():
    from extractor.chunker import chunk_text
    chunks = chunk_text("Hello world this is a short text.")
    assert len(chunks) == 1, "Short text should produce exactly 1 chunk"


# ── 5. Embeddings (LIVE — hits OpenAI) ──────────────────────────────────────
@pytest.mark.asyncio
async def test_get_embedding_returns_vector():
    from llm.openai_client import get_embedding
    vector = await get_embedding("Newton second law force mass acceleration")
    assert isinstance(vector, list), "Should return a list"
    assert len(vector) == 1536, f"text-embedding-3-small should return 1536 dims, got {len(vector)}"
    assert all(isinstance(v, float) for v in vector), "All values should be floats"
    print(f"\n[EMBEDDING] Vector dim={len(vector)}, first 3 values={vector[:3]}")


@pytest.mark.asyncio
async def test_get_embedding_empty_returns_empty():
    from llm.openai_client import get_embedding
    result = await get_embedding("")
    assert result == [], "Empty text should return empty list"


# ── 6. Pinecone upsert + search (LIVE — hits real Pinecone) ─────────────────
@pytest.mark.asyncio
async def test_pinecone_upsert_and_search():
    from rag.vector_store import upsert_vectors, search_vectors
    from llm.openai_client import get_embedding

    session_id = "test-session-ingest-001"
    document_id = "test-doc-001"
    filename = "test_physics.pdf"
    chunks = ["Newton second law states F equals ma. Force equals mass times acceleration."]

    # Embed
    vectors = [await get_embedding(chunks[0])]
    assert len(vectors) == 1 and len(vectors[0]) == 1536

    # Upsert
    stored = await upsert_vectors(session_id, document_id, filename, chunks, vectors)
    assert stored == 1, f"Expected 1 chunk stored, got {stored}"
    print(f"\n[PINECONE] Upserted {stored} vector(s)")

    # Wait briefly for Pinecone indexing
    await asyncio.sleep(2)

    # Search
    query_vec = await get_embedding("What is Newton second law?")
    results = await search_vectors(session_id, query_vec, top_k=5)
    assert len(results) > 0, "Should retrieve at least 1 result"
    assert results[0]["filename"] == filename, f"Filename mismatch: {results[0]['filename']}"
    assert "Newton" in results[0]["text"] or "Force" in results[0]["text"]
    print(f"\n[PINECONE] Retrieved {len(results)} result(s). Top score={results[0]['score']:.3f}")


# ── 7. OpenAI chat completion (LIVE — hits gpt-4o-mini directly) ────────────
@pytest.mark.asyncio
async def test_get_chat_response_returns_text():
    from llm.openai_client import get_chat_response
    messages = [
        {"role": "system", "content": "You are a helpful assistant. Answer in one sentence."},
        {"role": "user", "content": "What does F=ma mean?"}
    ]
    reply = await get_chat_response(messages)
    assert isinstance(reply, str), "Should return a string"
    assert len(reply) > 10, "Reply should not be empty"
    assert any(word in reply.lower() for word in ["force", "mass", "acceleration", "newton"]), (
        f"Expected physics content in reply, got: {reply}"
    )
    print(f"\n[GPT-4o-mini] Reply: {reply}")


# ── 8. Full /ingest endpoint (LIVE HTTP — server must be running on :8001) ───
SERVICE_KEY = "studdybuddy_internal_2026"
BASE_URL = "http://localhost:8001"


@pytest.mark.asyncio
async def test_ingest_endpoint_pdf():
    pdf_bytes = make_test_pdf_bytes()
    b64 = base64.b64encode(pdf_bytes).decode()

    payload = {
        "file_b64": b64,
        "filename": "test_physics.pdf",
        "mimetype": "application/pdf",
        "sessionId": "test-session-http-001",
        "documentId": "test-doc-http-001",
    }

    async with httpx.AsyncClient(timeout=60) as client:
        resp = await client.post(
            f"{BASE_URL}/ingest",
            json=payload,
            headers={"x-service-key": SERVICE_KEY}
        )

    assert resp.status_code == 200, f"Expected 200, got {resp.status_code}: {resp.text}"
    data = resp.json()
    assert data["success"] is True, f"success should be True: {data}"
    assert data["chunks_stored"] > 0, f"Should have stored at least 1 chunk: {data}"
    print(f"\n[HTTP /ingest] Status={resp.status_code}, chunks_stored={data['chunks_stored']}")


@pytest.mark.asyncio
async def test_ingest_endpoint_rejects_wrong_key():
    async with httpx.AsyncClient(timeout=10) as client:
        resp = await client.post(
            f"{BASE_URL}/ingest",
            json={"file_b64": "aaa", "filename": "x", "mimetype": "application/pdf",
                  "sessionId": "x", "documentId": "x"},
            headers={"x-service-key": "wrong-key"}
        )
    assert resp.status_code == 401, f"Expected 401 for wrong key, got {resp.status_code}"
    print(f"\n[HTTP /ingest] Wrong key correctly rejected with 401")
